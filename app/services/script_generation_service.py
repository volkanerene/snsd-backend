"""
Script Generation Service
Generates video scripts using ChatGPT for education topics, PDFs, and incident reports
"""

import json
import asyncio
import base64
from io import BytesIO
from typing import Optional, Dict, Any, List
from app.config import settings

try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def _extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF bytes"""
    if PyPDF2 is None:
        raise ValueError("PyPDF2 not installed. Install with: pip install PyPDF2")

    pdf_file = BytesIO(file_bytes)
    pdf_reader = PyPDF2.PdfReader(pdf_file)

    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"

    return text.strip()


def _extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from PPTX bytes"""
    if Presentation is None:
        raise ValueError("python-pptx not installed. Install with: pip install python-pptx")

    presentation = Presentation(BytesIO(file_bytes))
    texts: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts).strip()


def _parse_chat_message_content(content: Any) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict):
                if item.get("type") in {"text", "output_text"} and item.get("text"):
                    parts.append(item["text"])
        return "\n".join(parts)
    return ""


async def _extract_text_from_image(file_bytes: bytes, mime_type: str) -> str:
    """Use OpenAI vision to extract text from images"""
    if OpenAI is None or not settings.OPENAI_API_KEY:
        raise ValueError("OpenAI API key not configured for image extraction")

    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    encoded = base64.b64encode(file_bytes).decode("utf-8")
    data_url = f"data:{mime_type};base64,{encoded}"

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You read images and return the exact text content. Respond with plain text only."
            },
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all readable text from this document image."},
                    {"type": "image_url", "image_url": {"url": data_url}}
                ]
            }
        ],
        max_tokens=700,
    )

    content = response.choices[0].message.content
    extracted = _parse_chat_message_content(content).strip()
    if not extracted:
        raise ValueError("Could not extract text from image")
    return extracted


async def extract_text_from_document(
    filename: str,
    content_type: Optional[str],
    file_bytes: bytes
) -> str:
    """Determine document type and extract text accordingly"""
    lowered = (filename or "").lower()
    content_type = content_type or ""

    if lowered.endswith(".pdf") or content_type == "application/pdf":
        return _extract_text_from_pdf(file_bytes)

    if lowered.endswith((".pptx", ".ppt")) or content_type in {
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }:
        return _extract_text_from_pptx(file_bytes)

    if content_type.startswith("image/") or lowered.endswith((".png", ".jpg", ".jpeg", ".webp")):
        return await _extract_text_from_image(file_bytes, content_type or "image/png")

    raise ValueError(f"Unsupported file type for {filename or 'document'}")


def _build_education_prompt(topic: str) -> str:
    """Build prompt for education script generation"""
    return f"""You are a professional training video scriptwriter specializing in safety,
health, and business training content.

Create a compelling and engaging video script for the following educational topic:

TOPIC: {topic}

The script should:
1. Have a clear introduction that grabs attention
2. Present 3-5 key learning points with examples
3. Include practical takeaways
4. End with a memorable conclusion

IMPORTANT FORMAT RULES:
- Write the entire response as spoken narration or dialogue only.
- Do NOT include headings, bullet points, scene directions, or stage notes.
- No brackets, numbered lists, or meta commentary. Every line should sound like someone speaking naturally.

Keep sentences concise and engaging. Use a professional but friendly tone."""


def _build_pdf_script_prompt(pdf_content: str) -> str:
    """Build prompt for PDF-based script generation"""
    return f"""You are a professional training video scriptwriter. Your task is to convert
educational material into an engaging video script.

Here is the educational content from a PDF (trimmed for length):

---
{pdf_content[:3000]}
---

Create a video script that:
1. Extracts the most important information from this content
2. Presents it in an engaging, conversational way
3. Includes practical examples or applications
4. Has a clear beginning, middle, and end

STRICT FORMAT RULES:
- Output MUST be pure spoken dialogue/narration with complete sentences.
- Do NOT add headings, bullet points, numbered lists, scene directions, or stage notes.
- Avoid meta phrases like "In this script" or "Scene 1". Every line should sound like the narrator speaking directly to the audience.

Keep it suitable for a 2-3 minute video and write ONLY the spoken script."""


def _build_incident_script_prompt(
    what_happened: str,
    why_did_it_happen: Optional[str],
    what_did_they_learn: Optional[str],
    ask_yourself_or_crew: Optional[str],
    similar_incident: Optional[Dict[str, Any]] = None,
    process_safety_violations: Optional[str] = None,
    life_saving_rule_violations: Optional[str] = None,
    preventive_actions: Optional[str] = None,
    reference_case: Optional[str] = None
) -> str:
    """Build prompt for incident-based script generation with conservative, data-respecting approach.

    New approach:
    - Prevents data invention/hallucination
    - Conditionally includes sections only if data provided
    - Produces 350-450 word scripts
    - Single continuous spoken narrative
    - Graceful handling of missing data
    """

    # Build conditional reference case section
    reference_case_section = ""
    if similar_incident:
        ref_title = similar_incident.get('reference_case_title', similar_incident.get('title', ''))
        ref_desc = similar_incident.get('reference_case_description', '')
        ref_year = similar_incident.get('reference_case_year', '')
        if ref_title or ref_desc:
            reference_case_section = f"""Reference incident data:
Title: {ref_title}
Description: {ref_desc}
Year: {ref_year}

Include this in the script ONLY if you can clearly explain the specific mechanism similarity to the current incident."""

    # Build conditional PSF section
    psf_section = ""
    if process_safety_violations and process_safety_violations.strip():
        psf_section = f"""Process Safety Fundamentals violated: {process_safety_violations}
If PSF data is provided, briefly state which one and the specific mechanism that violated it."""

    # Build conditional LSR section
    lsr_section = ""
    if life_saving_rule_violations and life_saving_rule_violations.strip():
        lsr_section = f"""Life Saving Rules violated: {life_saving_rule_violations}
If LSR data is provided, briefly state which one and why it was not followed."""

    # Build conditional preventive actions section
    actions_section = ""
    if preventive_actions and preventive_actions.strip():
        actions_section = f"""Preventive actions documented: {preventive_actions}
If specific actions are provided, mention them naturally. If not, reference only general prevention principles."""

    return f"""You are Marcel, a safety expert delivering a 60–90 second learning-from-incident video.
You speak calmly, clearly, and professionally as if talking to frontline workers and supervisors.
Your output is a SINGLE continuous spoken script with NO headings, NO lists, NO bullet points, NO numbering.

CRITICAL RULES — NEVER BREAK THESE:
1. DO NOT invent any information that is not explicitly provided.
2. If "why it happened" is missing or unclear: Say it's "not fully detailed in the available information"
3. If "preventive actions" are not provided: Give only general prevention principles from the incident facts.
4. If no similar incident provided: COMPLETELY skip any industry reference case.
5. If PSF or LSR information is not provided: Do NOT guess or invent them—skip or say they cannot be identified from available data.
6. EVERYTHING must be traceable to incident data exactly as given.
7. Speak naturally with no lists, headings, or section labels.
8. Use the learnings/questions provided if available; if not, derive them carefully from the facts only.

INTERNAL STRUCTURE (do NOT show as sections in output):
- Introduction: 1–2 sentences engaging the listener
- What happened: Chronological facts from provided data only
- Why it happened: ONLY if provided in detail; otherwise state it's not fully documented
- Reference case: ONLY if similar incident provided AND you can explain specific mechanism similarity
- What can be learned: From provided learnings OR derived carefully from facts
- PSF/LSR statement: ONLY if specific data provided about violations
- Reflective question: Encourage thoughtful consideration

INCIDENT DATA PROVIDED:
What happened: {what_happened}
Why it happened: {why_did_it_happen if why_did_it_happen and why_did_it_happen.strip() else 'NOT PROVIDED'}
What we learned: {what_did_they_learn if what_did_they_learn and what_did_they_learn.strip() else 'NOT PROVIDED'}
Ask yourself or crew: {ask_yourself_or_crew if ask_yourself_or_crew and ask_yourself_or_crew.strip() else 'NOT PROVIDED'}

{reference_case_section}

{psf_section}

{lsr_section}

{actions_section}

FINAL CONSTRAINTS:
- Maximum: 350–450 words (~2200 characters)
- ABSOLUTELY NO invented data, numbers, equipment names, or mechanisms not explicitly in the provided incident information.
- If data is missing, state it naturally: "The specific cause isn't fully detailed in what we have," not invent reasons.
- Speak directly to the audience as if in conversation, not like reading a prepared report.
- Make it traceable: A reader should be able to point to the provided data and find the source of every claim you make."""


def _strip_json_block(raw_text: str) -> str:
    """Remove Markdown fences or stray text before attempting json.loads"""
    cleaned = raw_text.strip()
    if cleaned.startswith("```"):
        cleaned = cleaned.strip("`")
        if cleaned.startswith("json"):
            cleaned = cleaned[4:]
    if cleaned.endswith("```"):
        cleaned = cleaned[:-3]
    return cleaned.strip()


async def summarize_incident_from_text(document_text: str) -> Dict[str, str]:
    """
    Use GPT to summarize structured incident answers from extracted text.
    Returns dict with keys: what_happened, why_did_it_happen, what_did_they_learn, ask_yourself_or_crew.
    """
    if OpenAI is None or not settings.OPENAI_API_KEY:
        raise ValueError("OpenAI API key not configured")

    trimmed = (document_text or "").strip()
    if not trimmed:
        raise ValueError("No document text to analyze")

    trimmed = trimmed[:15000]

    prompt = f"""You are a safety analyst. Read the following incident report text and extract the requested
sections. If a section is missing, return an empty string for that field. Each summary must be 2-4
concise sentences (max 600 characters).

Respond ONLY with valid JSON (no markdown) in this format:
{{
  "what_happened": "...",
  "why_did_it_happen": "...",
  "what_did_they_learn": "...",
  "ask_yourself_or_crew": "A challenging reflection question derived from the document (or empty)."
}}

INCIDENT REPORT TEXT:
---
{trimmed}
---
"""

    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You extract structured safety-incident summaries and reply with strict JSON."
            },
            {"role": "user", "content": prompt}
        ],
        temperature=0.2,
        max_tokens=700
    )

    raw_content = response.choices[0].message.content or ""
    cleaned = _strip_json_block(raw_content)

    try:
        parsed = json.loads(cleaned)
    except json.JSONDecodeError:
        raise ValueError("Failed to parse incident summary JSON from model response")

    return {
        "what_happened": (parsed.get("what_happened") or "").strip(),
        "why_did_it_happen": (parsed.get("why_did_it_happen") or "").strip(),
        "what_did_they_learn": (parsed.get("what_did_they_learn") or "").strip(),
        "ask_yourself_or_crew": (parsed.get("ask_yourself_or_crew") or "").strip(),
    }


async def generate_script_from_topic(topic: str) -> Dict[str, Any]:
    """Generate a video script from an educational topic"""
    try:
        if not settings.OPENAI_API_KEY:
            return {
                "success": False,
                "error": "OpenAI API key not configured"
            }

        if OpenAI is None:
            return {
                "success": False,
                "error": "OpenAI library not installed"
            }

        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        prompt = _build_education_prompt(topic)

        print(f"[Script Gen] Generating script from topic: {topic[:50]}...")

        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {
                    "role": "system",
                    "content": "You are a professional training video scriptwriter."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=2000
        )

        script = response.choices[0].message.content.strip()

        return {
            "success": True,
            "script": script,
            "source": "topic",
            "generated_at": str(__import__('datetime').datetime.now(
                __import__('datetime').timezone.utc
            ).isoformat())
        }

    except Exception as e:
        print(f"[Script Gen] Error: {str(e)}")
        return {
            "success": False,
            "error": f"Script generation failed: {str(e)}"
        }


async def generate_script_from_material(material: str) -> Dict[str, Any]:
    """Generate a video script from textual training material"""
    try:
        if not settings.OPENAI_API_KEY:
            return {
                "success": False,
                "error": "OpenAI API key not configured"
            }

        if OpenAI is None:
            return {
                "success": False,
                "error": "OpenAI library not installed"
            }

        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        prompt = _build_pdf_script_prompt(material)

        print(f"[Script Gen] Generating script from material ({len(material)} chars)...")

        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {
                    "role": "system",
                    "content": "You are a professional training video scriptwriter."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=2000
        )

        script = response.choices[0].message.content.strip()

        return {
            "success": True,
            "script": script,
            "source": "document",
            "generated_at": str(__import__('datetime').datetime.now(
                __import__('datetime').timezone.utc
            ).isoformat())
        }

    except Exception as e:
        print(f"[Script Gen] Error: {str(e)}")
        return {
            "success": False,
            "error": f"Document script generation failed: {str(e)}"
        }


# Backwards compatibility alias
generate_script_from_pdf = generate_script_from_material


def _generate_training_questions(
    what_happened: str,
    why_did_it_happen: Optional[str],
    what_did_they_learn: Optional[str]
) -> Dict[str, Any]:
    """Generate 3 training questions: 2 multiple choice, 1 text-based"""
    try:
        if OpenAI is None or not settings.OPENAI_API_KEY:
            return {"success": False, "questions": []}

        client = OpenAI(api_key=settings.OPENAI_API_KEY)

        prompt = f"""Based on this incident, generate exactly 3 training questions.

INCIDENT:
What happened: {what_happened}
Why it happened: {why_did_it_happen or 'N/A'}
What we learned: {what_did_they_learn or 'N/A'}

Generate exactly this JSON format (no markdown, pure JSON):
{{
  "questions": [
    {{
      "type": "multiple_choice",
      "question": "question text here?",
      "options": ["option A", "option B", "option C", "option D"],
      "correct_answer": 0
    }},
    {{
      "type": "multiple_choice",
      "question": "another question?",
      "options": ["option A", "option B", "option C", "option D"],
      "correct_answer": 1
    }},
    {{
      "type": "text",
      "question": "What would you do differently in this situation?"
    }}
  ]
}}

Keep questions professional and related to the incident. Make them assessment-focused."""

        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {
                    "role": "system",
                    "content": "You are a training assessment expert. Generate only valid JSON with no additional text."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=800
        )

        response_text = response.choices[0].message.content.strip()

        # Parse JSON response
        questions_data = json.loads(response_text)
        return {
            "success": True,
            "questions": questions_data.get("questions", [])
        }

    except Exception as e:
        print(f"[Training Questions] Error: {str(e)}")
        return {"success": False, "questions": []}


def generate_questions_from_script(script_text: str) -> Dict[str, Any]:
    """
    Generate two multiple-choice and one short-answer question
    based on the provided training script.
    """
    try:
        if OpenAI is None or not settings.OPENAI_API_KEY:
            return {"success": False, "questions": []}

        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        trimmed_script = script_text[:4000]

        prompt = f"""You are a training assessment expert.
Read the following training script and produce assessment questions to verify comprehension.

SCRIPT:
---
{trimmed_script}
---

Respond with valid JSON (no markdown) using this structure:
{{
  "questions": [
    {{
      "type": "multiple_choice",
      "question": "...",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_answer": 0
    }},
    {{
      "type": "multiple_choice",
      "question": "...",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_answer": 2
    }},
    {{
      "type": "text",
      "question": "Short-answer reflection question"
    }}
  ]
}}

Guidelines:
- Both multiple choice questions must have exactly 4 options with zero-based correct_answer index.
- The final question must be open-ended (type = "text").
- Focus on practical safety takeaways referenced in the script."""

        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {
                    "role": "system",
                    "content": "You generate structured assessment questions and only return strict JSON output."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.4,
            max_tokens=700
        )

        response_text = response.choices[0].message.content.strip()
        questions_data = json.loads(response_text)
        return {
            "success": True,
            "questions": questions_data.get("questions", [])
        }

    except Exception as e:
        print(f"[Training Questions] Error generating from script: {str(e)}")
        return {"success": False, "questions": []}


async def generate_script_from_incident(
    what_happened: str,
    why_did_it_happen: Optional[str] = None,
    what_did_they_learn: Optional[str] = None,
    ask_yourself_or_crew: Optional[str] = None,
    tenant_id: Optional[str] = None,
    process_safety_violations: Optional[str] = None,
    life_saving_rule_violations: Optional[str] = None,
    preventive_actions: Optional[str] = None,
    reference_case: Optional[str] = None
) -> Dict[str, Any]:
    """
    Generate a video script from incident details with conservative, data-respecting approach.
    Optionally finds similar incidents from database for context.

    Key features:
    - Prevents hallucination by never inventing data
    - Handles missing data gracefully (states when not detailed rather than inventing)
    - Produces 350-450 word single continuous spoken scripts
    - Conditionally includes sections (PSF, LSR, reference cases) only if data provided
    - Makes every claim traceable to provided incident data
    """
    try:
        if not settings.OPENAI_API_KEY:
            return {
                "success": False,
                "error": "OpenAI API key not configured"
            }

        if OpenAI is None:
            return {
                "success": False,
                "error": "OpenAI library not installed"
            }

        # Find similar incident from database if tenant_id provided
        similar_incident = None
        if tenant_id:
            similar_incident = _find_similar_incident(
                what_happened,
                tenant_id
            )
            if similar_incident:
                print(f"[Script Gen] Found similar incident: {similar_incident.get('title', 'N/A')}")

        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        prompt = _build_incident_script_prompt(
            what_happened,
            why_did_it_happen,
            what_did_they_learn,
            ask_yourself_or_crew,
            similar_incident,
            process_safety_violations,
            life_saving_rule_violations,
            preventive_actions,
            reference_case
        )

        print(f"[Script Gen] Generating incident script with conservative, data-respecting approach...")

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are Marcel, a safety expert. Create natural, conversational dialogue only - no scene descriptions, no stage directions, no visual cues, no headings, no lists. Output a single continuous spoken script for a 60-90 second video. Never invent data—only use information explicitly provided. If data is missing, state it naturally."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.6,
            max_tokens=1500
        )

        script = response.choices[0].message.content.strip()

        # Generate training questions
        questions_result = _generate_training_questions(
            what_happened,
            why_did_it_happen,
            what_did_they_learn
        )

        return {
            "success": True,
            "script": script,
            "source": "incident",
            "questions": questions_result.get("questions", []),
            "similar_incident_id": similar_incident.get("id") if similar_incident else None,
            "generated_at": str(__import__('datetime').datetime.now(
                __import__('datetime').timezone.utc
            ).isoformat())
        }

    except Exception as e:
        print(f"[Script Gen] Error: {str(e)}")
        return {
            "success": False,
            "error": f"Incident script generation failed: {str(e)}"
        }


def _find_similar_incident(
    description: str,
    tenant_id: str
) -> Optional[Dict[str, Any]]:
    """
    Find the most similar incident from database using semantic similarity.
    Searches both tenant-specific incidents and global templates.
    For now, using simple keyword matching. Can be enhanced with embeddings later.
    """
    try:
        from app.db.supabase_client import supabase

        # Get incidents for tenant + global templates
        # First get tenant-specific incidents
        res_tenant = supabase.table("incident_report_dialogues").select(
            "id, title, what_happened, why_did_it_happen, what_did_they_learn"
        ).eq("tenant_id", tenant_id).execute()

        # Then get global template incidents (is_template=true, tenant_id=NULL)
        res_templates = supabase.table("incident_report_dialogues").select(
            "id, title, what_happened, why_did_it_happen, what_did_they_learn"
        ).eq("is_template", True).is_("tenant_id", "null").execute()

        incidents = []
        if hasattr(res_tenant, 'data') and res_tenant.data:
            incidents.extend(res_tenant.data)
        if hasattr(res_templates, 'data') and res_templates.data:
            incidents.extend(res_templates.data)

        if not incidents:
            return None

        # Simple keyword matching: count common words
        description_words = set(description.lower().split())

        best_match = None
        best_score = 0

        for incident in incidents:
            incident_text = (
                incident.get("what_happened", "").lower() + " " +
                incident.get("why_did_it_happen", "").lower()
            )
            incident_words = set(incident_text.split())

            # Calculate Jaccard similarity
            if incident_words and description_words:
                intersection = len(description_words & incident_words)
                score = intersection / len(description_words | incident_words)

                if score > best_score:
                    best_score = score
                    best_match = incident

        # Only return if similarity > 0.1 (at least some common words)
        return best_match if best_score > 0.1 else None

    except Exception as e:
        print(f"[Script Gen] Error finding similar incident: {str(e)}")
        return None
